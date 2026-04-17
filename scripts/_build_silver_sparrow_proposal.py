"""
Build Valuation Realized Commercial Proposal - Canada (Silver Sparrow)
Mirrors the UAE proposal structure but re-scoped for Pramod Kumar's fundraise,
with 4 engagement buckets and CAD pricing.
FX: USD quotes from the call translated at 1.40 USD/CAD (indicative).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# Palette (extracted from the UAE template visuals)
GREEN = RGBColor(0x15, 0x32, 0x23)          # dark forest green
GOLD = RGBColor(0xC6, 0x9A, 0x55)            # gold/bronze accent
CREAM = RGBColor(0xF2, 0xED, 0xE3)           # cream page background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEA, 0xE6, 0xDC)           # alternating row shading
MUTED = RGBColor(0x4A, 0x4A, 0x4A)           # muted body text
SIDE_TEXT = RGBColor(0xD9, 0xCF, 0xB8)       # light cream text on dark sidebar

HEADER_FONT = "Georgia"
BODY_FONT = "Calibri"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_text(slide, left, top, width, height, text, *, font=BODY_FONT, size=14,
             bold=False, italic=False, color=MUTED, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_gold_top_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.13))
    set_fill(bar, GOLD)


def add_rect(slide, left, top, width, height, fill_rgb, *, line=False):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_rgb
    if line:
        rect.line.color.rgb = fill_rgb
    else:
        rect.line.fill.background()
    return rect


# ---------------- Build ----------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


# ===== Slide 1: Title =====
s1 = prs.slides.add_slide(blank_layout)
bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
set_fill(bg, GREEN)
add_gold_top_bar(s1)

add_text(s1, Inches(0.8), Inches(2.3), Inches(11.5), Inches(1.8),
         "Valuation Realized – Canada",
         font=HEADER_FONT, size=60, bold=True, color=WHITE)
add_text(s1, Inches(0.8), Inches(3.9), Inches(11.5), Inches(0.6),
         "Advisory Services",
         font=BODY_FONT, size=26, color=GOLD)
# gold underline
und = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.7),
                          Inches(1.3), Inches(0.04))
set_fill(und, GOLD)
add_text(s1, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.4),
         "Non-binding scoping proposal – Silver Sparrow (software arm)",
         font=BODY_FONT, size=14, color=SIDE_TEXT)
add_text(s1, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
         "Confidential",
         font=BODY_FONT, size=12, italic=True, color=SIDE_TEXT)


# ===== Slide 2: Our approach =====
s2 = prs.slides.add_slide(blank_layout)
bg = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
set_fill(bg, CREAM)
add_gold_top_bar(s2)

add_text(s2, Inches(0.7), Inches(0.55), Inches(12), Inches(1.0),
         "Our approach",
         font=HEADER_FONT, size=44, bold=True, color=GREEN)
add_text(s2, Inches(0.7), Inches(1.55), Inches(12), Inches(1.2),
         ("Valuation Realized provides independent M&A and capital-raise advisory to "
          "owner-managed businesses. Based on our call, the services below are "
          "structured across four engagement tiers, each available standalone or "
          "combined. All fees are quoted in Canadian dollars."),
         font=BODY_FONT, size=16, color=GREEN)

# Four cards
cards = [
    ("01", "Business\nvaluation", "Defensible, market-calibrated number to anchor investor conversations"),
    ("02", "Investor\npack", "Credible pitch deck and teaser to substantiate the story and raise"),
    ("03", "Sale or raise\nmandate", "End-to-end support from outreach through to close"),
    ("04", "Hourly\nadvisory", "Targeted inputs, ad-hoc reviews, no minimum commitment"),
]
card_w = Inches(2.95)
card_h = Inches(3.0)
card_y = Inches(3.4)
gap = Inches(0.15)
start_x = Inches(0.7)
for i, (num, title, desc) in enumerate(cards):
    x = start_x + (card_w + gap) * i
    c = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, card_y, card_w, card_h)
    set_fill(c, WHITE)
    add_text(s2, x + Inches(0.25), card_y + Inches(0.25), card_w - Inches(0.5),
             Inches(0.7), num,
             font=HEADER_FONT, size=34, bold=True, color=GOLD)
    add_text(s2, x + Inches(0.25), card_y + Inches(0.95), card_w - Inches(0.5),
             Inches(1.1), title,
             font=HEADER_FONT, size=22, bold=True, color=GREEN)
    add_text(s2, x + Inches(0.25), card_y + Inches(2.05), card_w - Inches(0.5),
             Inches(0.85), desc,
             font=BODY_FONT, size=12, color=MUTED)


# ---- helper to build the detail slides ----
def add_detail_slide(num, title, subtitle, items, price_block):
    s = prs.slides.add_slide(blank_layout)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_fill(bg, CREAM)
    add_gold_top_bar(s)

    # Title: "01  Business valuation"
    add_text(s, Inches(0.7), Inches(0.55), Inches(8.5), Inches(1.0),
             f"{num}   {title}",
             font=HEADER_FONT, size=40, bold=True, color=GREEN)

    # subtitle
    add_text(s, Inches(0.7), Inches(1.55), Inches(7.8), Inches(1.0),
             subtitle,
             font=BODY_FONT, size=14, color=GREEN)

    # bullets as alternating rows
    row_h = Inches(0.85)
    row_w = Inches(7.9)
    row_y = Inches(2.85)
    for i, (head, sub) in enumerate(items):
        y = row_y + row_h * i
        bg_rgb = WHITE if i % 2 == 0 else LIGHT
        row = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), y, row_w, row_h)
        set_fill(row, bg_rgb)
        add_text(s, Inches(0.9), y + Inches(0.1), row_w - Inches(0.4),
                 Inches(0.35), head,
                 font=BODY_FONT, size=13, bold=True, color=GREEN)
        add_text(s, Inches(0.9), y + Inches(0.45), row_w - Inches(0.4),
                 Inches(0.4), sub,
                 font=BODY_FONT, size=11, color=MUTED)

    # pricing sidebar
    side_x = Inches(9.0)
    side_y = Inches(1.55)
    side_w = Inches(3.6)
    side_h = Inches(5.2)
    side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, side_x, side_y, side_w, side_h)
    set_fill(side, GREEN)

    add_text(s, side_x + Inches(0.3), side_y + Inches(0.3),
             side_w - Inches(0.6), Inches(0.4), "Pricing",
             font=BODY_FONT, size=16, color=GOLD)

    # Fill the price block lines
    cur_y = side_y + Inches(0.85)
    for item in price_block:
        kind = item["kind"]
        text = item["text"]
        if kind == "label":
            add_text(s, side_x + Inches(0.3), cur_y, side_w - Inches(0.6),
                     Inches(0.4), text,
                     font=BODY_FONT, size=14, bold=True, color=WHITE)
            cur_y += Inches(0.5)
        elif kind == "big":
            add_text(s, side_x + Inches(0.3), cur_y, side_w - Inches(0.6),
                     Inches(0.9), text,
                     font=HEADER_FONT, size=34, bold=True, color=GOLD)
            cur_y += Inches(0.95)
        elif kind == "sub":
            add_text(s, side_x + Inches(0.3), cur_y, side_w - Inches(0.6),
                     Inches(0.4), text,
                     font=BODY_FONT, size=12, color=SIDE_TEXT)
            cur_y += Inches(0.45)
        elif kind == "rule":
            rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      side_x + Inches(0.3), cur_y + Inches(0.15),
                                      side_w - Inches(0.6), Emu(12700))
            set_fill(rule, GOLD)
            cur_y += Inches(0.35)
        elif kind == "spacer":
            cur_y += Inches(float(text))
        elif kind == "note":
            add_text(s, side_x + Inches(0.3), cur_y, side_w - Inches(0.6),
                     Inches(0.4), text,
                     font=BODY_FONT, size=11, color=SIDE_TEXT)
            cur_y += Inches(0.4)


# ===== Slide 3: 01 Business valuation =====
add_detail_slide(
    "01",
    "Business valuation",
    ("Independent valuation providing a defensible, market-calibrated number to "
     "anchor investor conversations. Pre-revenue methodologies applied, given "
     "Silver Sparrow's software arm is carve-out and pre-revenue."),
    [
        ("Venture capital / pre-revenue method",
         "Target return-based valuation calibrated to raise stage"),
        ("Comparable transaction benchmarking",
         "TMS and logistics SaaS comparables (Canadian and international)"),
        ("Projection review & sanity check",
         "Revenue build, unit economics, capital efficiency"),
        ("Scenario & sensitivity",
         "Best/worst/base case tied to adoption and sales ramp assumptions"),
    ],
    [
        {"kind": "label", "text": "Fixed fee"},
        {"kind": "big", "text": "CAD 2,400"},
        {"kind": "sub", "text": "Indicative, based on current scope"},
        {"kind": "rule", "text": ""},
        {"kind": "sub", "text": "50% on engagement"},
        {"kind": "sub", "text": "50% on delivery"},
    ],
)

# ===== Slide 4: 02 Investor pack =====
add_detail_slide(
    "02",
    "Investor pack",
    ("Prepare the investor-facing documentation pack. Build a credible story "
     "that substantiates the valuation and filters qualified investors faster."),
    [
        ("Pitch deck refresh / rebuild",
         "Starting from your existing deck; story, structure, and visual polish"),
        ("Confidential teaser",
         "1-page blind summary for initial outreach"),
        ("Financial projections review",
         "Link to valuation, assumptions, and use of proceeds"),
        ("NDA + process framework",
         "Legal templates and investor qualification workflow"),
    ],
    [
        {"kind": "label", "text": "Fixed fee"},
        {"kind": "big", "text": "CAD 3,200"},
        {"kind": "sub", "text": "Indicative, based on current scope"},
        {"kind": "rule", "text": ""},
        {"kind": "sub", "text": "50% on engagement"},
        {"kind": "sub", "text": "50% on delivery of pack"},
    ],
)

# ===== Slide 5: 03 Sale or raise mandate =====
add_detail_slide(
    "03",
    "Sale or raise mandate",
    ("End-to-end support for your capital raise or sale process, from targeted "
     "outreach through negotiation to close. Structured as a low retainer + "
     "success fee on outcome."),
    [
        ("Targeted investor / buyer outreach",
         "Strategic buyers, VCs, logistics operators, industry contacts"),
        ("Qualification & NDA management",
         "Screening, scoring, confidentiality process"),
        ("Data room preparation & Q&A",
         "Structured diligence facilitation"),
        ("Negotiation support & deal structuring",
         "Price, earn-outs, liquidation preferences, transition terms"),
        ("Term sheet / SPA review & close",
         "Through to signed agreement and handover"),
    ],
    [
        {"kind": "label", "text": "Retainer"},
        {"kind": "big", "text": "CAD 700"},
        {"kind": "sub", "text": "per week"},
        {"kind": "rule", "text": ""},
        {"kind": "label", "text": "Success fee"},
        {"kind": "sub", "text": "7% on funds raised / total consideration ≤ CAD 2M"},
        {"kind": "sub", "text": "4% above CAD 2M"},
    ],
)

# ===== Slide 6: 04 Hourly advisory =====
add_detail_slide(
    "04",
    "Hourly advisory",
    ("For targeted inputs where a full engagement is not required, e.g. reviewing "
     "your existing deck, financial projections, or investor strategy. No minimum "
     "commitment."),
    [
        ("Deck & projections review",
         "Structured written feedback and recommended edits"),
        ("Investor strategy consultation",
         "Right investor profile, outreach approach, process design"),
        ("Ad-hoc working sessions",
         "Book time as needed; capped estimate agreed upfront per task"),
        ("Transition to fixed or retainer",
         "Hours credited if scope evolves into engagement tiers 01–03"),
    ],
    [
        {"kind": "label", "text": "Hourly rate"},
        {"kind": "big", "text": "CAD 310"},
        {"kind": "sub", "text": "per hour"},
        {"kind": "rule", "text": ""},
        {"kind": "sub", "text": "No minimum commitment"},
        {"kind": "sub", "text": "Billed monthly against time log"},
    ],
)


# ===== Slide 7: Thank you =====
s7 = prs.slides.add_slide(blank_layout)
bg = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
set_fill(bg, GREEN)
add_gold_top_bar(s7)

add_text(s7, Inches(0.8), Inches(1.8), Inches(11), Inches(1.6),
         "Thank you",
         font=HEADER_FONT, size=72, bold=True, color=WHITE)

und = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.6),
                          Inches(1.3), Inches(0.04))
set_fill(und, GOLD)

add_text(s7, Inches(0.8), Inches(4.1), Inches(11), Inches(0.5),
         "vr@valuationrealized.com",
         font=BODY_FONT, size=16, color=GOLD)
add_text(s7, Inches(0.8), Inches(4.6), Inches(11), Inches(0.5),
         "+971 505 832 701",
         font=BODY_FONT, size=16, color=GOLD)
add_text(s7, Inches(0.8), Inches(5.1), Inches(11), Inches(0.5),
         "valuationrealized.com",
         font=BODY_FONT, size=16, color=GOLD)

add_text(s7, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
         ("FX note: USD quotes from our call converted at indicative "
          "1.00 USD = 1.40 CAD. Final CAD invoicing will use the spot rate at "
          "engagement date."),
         font=BODY_FONT, size=10, italic=True, color=SIDE_TEXT)


# ---- save ----
out = r"C:\Users\vrimsaite\Desktop\Valuation Realized Commercial Proposal - Silver Sparrow (CAD).pptx"
prs.save(out)
print("WROTE:", out)
